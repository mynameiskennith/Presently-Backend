# import streamlit as st
from fastapi.responses import JSONResponse
import pptx
from pptx.util import Inches, Pt
import os
# from dotenv import load_dotenv
from groq import Groq
from io import BytesIO
import tempfile
import zipfile
import xml.etree.ElementTree as ET
# import time
from PIL import Image as Imagee
# import sounddevice as sd
# import numpy as np
# import scipy.io.wavfile as wavfile
# import speech_recognition as sr
from spire.presentation import Presentation as PPresentation
from spire.presentation.common import *
# import wave
# import struct
import json
import traceback
from io import BytesIO
from pptx import Presentation
import io, base64



# Get API key from environment variable
api_key = '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'
if not api_key:
    raise ValueError("GROQ_API_KEY not found in environment variables")

client = Groq(api_key=api_key)


# Custom Formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def generate_slide_titles(topic,noOfSlides,audienceType,slideContent):
    prompt = f"""Generate exactly {noOfSlides} concise slide titles for a presentation on the topic: {topic}
    Rules:
    1. Provide only the titles, one per line
    2. Do not include any numbering or bullet points
    3. Each title should be brief and relevant to the topic
    4. Do not include any additional text in response  or explanations
    6. Directly give the titles and do not add any additional message above it
    7. The type of presentation is {audienceType} 
    8. Give more importance to : {slideContent} when selecting topics """
    
    response = client.chat.completions.create(
        model="llama3-70b-8192",
        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ],
        temperature=0.7,
        max_tokens=1024,
        top_p=1,
        stream=False,
    )

    # Extract the content from the response
    response_text = response.choices[0].message.content
    # Split the response into titles and filter out empty lines
    return [title.strip() for title in response_text.split("\n") if title.strip()] 



def generate_slide_content(slide_title,audienceType):
    prompt = f"""Generate exactly 7 bullet points for the slide titled: "{slide_title}"
    Rules:
    1. Each point must be a very short but crisp sentence
    2. Do not exceed 15 words per point
    3. Provide only the points, one per line
    4. Do not include any numbering or bullet point symbols
    5. Do not include any additional text from response or 
    6. Each point should be self explanatory
    7. Directly provide the points for the slide title and do not include any additional message before the points
    8. Do not include the slide title in the points
    9. The type of presentation is {audienceType} , select the tone of points accordingly"""
    
    response = client.chat.completions.create(
        model="llama3-70b-8192",
        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ],
        temperature=0.7,
        max_tokens=1024,
        top_p=1,
        stream=False,
    )

    # Extract the content from the response
    response_text = response.choices[0].message.content

    # Split the response into points and filter out empty lines
    points = [point.strip() for point in response_text.split("\n") if point.strip()][1:7]  # Ensure we get exactly 6 points
    
    # Join the points with newlines to create the slide content
    return "\n".join(points)

def get_body_placeholder(slide):
    """Find the placeholder of type BODY (2) in a slide"""
    for shape in slide.shapes.placeholders:
        # Type 2 corresponds to BODY as seen in the debug output
        if shape.placeholder_format.type == 2:
            return shape
    return None

def create_presentation(request_data):
    template_path = f"templates/{request_data['template']}.pptx"
    print(template_path)
    if not os.path.exists(template_path):
        raise ValueError(f"Template {request_data['template']} not found")

    prs = Presentation(template_path)
    
    print("DEBUG: Analyzing template structure")
    for layout in prs.slide_layouts:
        print(f"Layout index: {prs.slide_layouts.index(layout)}")
        for placeholder in layout.placeholders:
            print(f"  Placeholder index: {placeholder.placeholder_format.idx}, type: {placeholder.placeholder_format.type}")
    print("DEBUG: Template analysis complete")

    # Add title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = request_data["topic"]

    # Add contents slide
    contents_slide = prs.slides.add_slide(prs.slide_layouts[1])
    contents_slide.shapes.title.text = "Contents"
    
    slide_titles = generate_slide_titles(
        request_data["topic"], 
        request_data["numberOfSlides"], 
        request_data["audienceType"], 
        request_data["slideContent"]
    )
    content_text = "\n".join(slide_titles)
    
    # Add content to contents slide
    body_placeholder = get_body_placeholder(contents_slide)
    if body_placeholder:
        body_placeholder.text = content_text
    
    # Add content slides
    slide_layout = prs.slide_layouts[2]  # Content slide layout
    for slide_title in slide_titles:
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide_content = generate_slide_content(slide_title, request_data["audienceType"])
        
        # Add content to slide
        body_placeholder = get_body_placeholder(slide)
        if body_placeholder:
            body_placeholder.text = slide_content

    # Add thank you slide
    thank_you_slide = prs.slides.add_slide(prs.slide_layouts[3])
    thank_you_slide.shapes.title.text = "Thank You"

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

############################################################################################################################################

# Function to extract all text from the PPT
def extract_text_from_ppt(ppt_file_contents):
    """Extract all text from the PPT for processing."""
    presentation = Presentation(BytesIO(ppt_file_contents))
    all_text = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        all_text.append("\n".join(slide_text))
    return "\n\n".join(all_text)

# Function to count the number of images in the PPT
def count_images_in_ppt(ppt_file_contents):
    """Count the number of images in the PPT."""
    presentation = Presentation(BytesIO(ppt_file_contents))
    total_images = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Shape type 13 corresponds to Picture
                total_images += 1
    return total_images



# Main function to rate the PPT
def rate_ppt(ppt_file_contents):
    """Evaluate a PPT file using the Groq model."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
        tmp_file.write(ppt_file_contents)
        tmp_file_path = tmp_file.name

    try:
        # Extract text and count images in the PPT
        ppt_text = extract_text_from_ppt(ppt_file_contents)
        num_images = count_images_in_ppt(ppt_file_contents)

        # Prepare the prompt for the Groq model
        prompt = f"""
        You are an AI assistant tasked with evaluating PowerPoint presentations.
        
        PPT content: {ppt_text}
        Number of images: {num_images}
        
        Analyze the provided presentation based on these detailed criteria of a professional presentation:
        
        1. **Number of points per slide**: 
            - Evaluate based on the 7/7, 8/8, or 6/6 rule (no more than 6-8 words per line and 6-8 lines per slide).
            - Deduct points if a slide exceeds these limits.
            - Score out of 100: Full points if all slides adhere, deductions for excess.
        
        2. **Number of images per slide**:
            - Number of Images = {num_images}.
            - Score = ( Number of Images/ Number of Slides ) * 100
            - Score out of 100: Full points for the ideal range, deductions for too few or too many images.
        
        3. **Readability of text content**:
            - Use the SMOG Readability Formula to assess readability (based on syllables and sentence complexity).
            - Full points for readability suited for the target audience (e.g., 6th-9th-grade level for general audiences).
            - Score out of 100.
        
        4. **Consistency of slide formatting**:
            - Check for uniformity in font styles, font sizes, color schemes, and alignment.
            - Deduct points for inconsistent elements across slides.
            - Score out of 100.
        
        5. **Overall presentation content quality**:
            - Assess the presentation's organization, logical flow, and coverage of the topic.
            - Deduct points for missing key information or lack of structure.
            - Score out of 100.
        
        6. **Number of slides**:
            - Check if the total number of slides is appropriate for the presentation's purpose (e.g., 8-12 slides for a 10-minute talk).
            - Deduct points for excessive or insufficient slides.
            - Score out of 100.
        
        7. **Overall score**:
            - An average based on all criteria above.
            - Score out of 100.

        Return a JSON object in this format:
        {{
            "noOfPoints": {{
                "score": <score out of 100>,
                "reason": "Detail reason for the score and how we computed the score"
            }},
            "noOfImages": {{
                "score": <score out of 100>,
                "reason": "Detail reason for the score and how we computed the score."
            }},
            "readability": {{
                "score": <score out of 100>,
                "reason": "Detail reason for the score and how we computed the score"
            }},
            "consistency": {{
                "score": <score out of 100>,
                "reason": "Detail reason for the score and how we computed the score"
            }},
            "quality": {{
                "score": <score out of 100>,
                "reason": "Detail reason for the score and how we computed the score"
            }},
            "noOfSlides": {{
                "score": <score out of 100>,
                "reason": "Detail reason for the score and how we computed the score"
            }},
            "overallScore": {{
                "score": <score out of 100>,
                "reason": "Detail reason for the score and how we computed the score"
            }}
        }}

        Only return the JSON object with no additional text.
        Strictly Follow this return Format at all cost 
        """

        # Pass the prompt to the Groq model
        response = client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            max_tokens=2048,
            top_p=1,
            stream=False,
        )

        # Extract and parse the response
        evaluation_results = response.choices[0].message.content.strip()
        # print(evaluation_results)
        try:
            evaluation_results = json.loads(evaluation_results)
            # print(evaluation_results)
        except json.JSONDecodeError:
            raise ValueError("Invalid JSON response from Groq API")

        return evaluation_results

    except Exception as error:
        print(traceback.format_exc())  # Log detailed error
        return {"Error": str(error)}

    finally:
        if os.path.exists(tmp_file_path):
            os.unlink(tmp_file_path)

##########################################################################################################

def transcribe_audio(audio_file_contents):
    with BytesIO(audio_file_contents) as audio_file:
        transcription = client.audio.transcriptions.create(
            file=("audio.m4a", audio_file.read()),
            model="whisper-large-v3",
            response_format="verbose_json",
        )
        print(transcription.text)
        return transcription.text

def analyse_training(ppt_file_contents, audio_file_contents, audio_length):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
        tmp_file.write(ppt_file_contents)
        tmp_file_path = tmp_file.name

    try:
        ppt_text = extract_text_from_ppt(ppt_file_contents)
        num_images = count_images_in_ppt(ppt_file_contents)

        transcribed = transcribe_audio(audio_file_contents)

        prompt = f"""
            You are an AI assistant tasked with evaluating PowerPoint presentations and transcribed audio.
            
            PPT content: {ppt_text}
            Number of images: {num_images}
            Audio Length : {audio_length}
            Transcribed Audio: {transcribed}
            
            Analyze the provided presentation based on these detailed criteria of a professional presentation, also make sure that
            the scores for delivery analysis clearly reflects the audio length and number of slides as users tend to test the voice recording
            by simply saying filler words for short amount of time and still the delivery analysis provides good scores. Make sure if the 
            audio length is short when compared to the slide number the scores are low :
            
            1. **Number of points per slide**: 
                - Evaluate based on the 7/7, 8/8, or 6/6 rule (no more than 6-8 words per line and 6-8 lines per slide).
                - Deduct points if a slide exceeds these limits.
                - Score out of 100: Full points if all slides adhere, deductions for excess.
            
            2. **Number of images per slide**:
                - Number of Images = {num_images}.
                - Score = ( Number of Images / Number of Slides ) * 100
                - Score out of 100: Full points for the ideal range, deductions for too few or too many images.
            
            3. **Readability of text content**:
                - Use the SMOG Readability Formula to assess readability (based on syllables and sentence complexity).
                - Full points for readability suited for the target audience (e.g., 6th-9th-grade level for general audiences).
                - Score out of 100.
            
            4. **Consistency of slide formatting**:
                - Check for uniformity in font styles, font sizes, color schemes, and alignment.
                - Deduct points for inconsistent elements across slides.
                - Score out of 100.
            
            5. **Overall presentation content quality**:
                - Assess the presentation's organization, logical flow, and coverage of the topic.
                - Deduct points for missing key information or lack of structure.
                - Score out of 100.
            
            6. **Number of slides**:
                - Check if the total number of slides is appropriate for the presentation's purpose (e.g., 8-12 slides for a 10-minute talk).
                - Deduct points for excessive or insufficient slides.
                - Score out of 100.
            
            7. **Audio delivery evaluation**:
                - Analyze the transcription of the audio for clarity, pace, engagement, confidence, and alignment with the slides.
                - Score based on factors such as clarity of speech, engagement, use of filler words, pace, and alignment with slide content.
            
            8. **Overall score**:
                - An average based on all criteria above.
                - Score out of 100.

            Return a JSON object in this format:
            {{
                "contentAnalysis": {{
                    "noOfPoints": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "noOfImages": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score."
                    }},
                    "readability": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "consistency": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "quality": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "noOfSlides": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "overallScore": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }}
                }},
                "deliveryAnalysis": {{
                    "clarity": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "pace": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "engagement": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "confidence": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "alignment": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "structure": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }},
                    "overallScore": {{
                        "score": <score out of 100>,
                        "reason": "Detail reason for the score and how we computed the score"
                    }}
                }},
                "combinedEvaluation": {{
                    "combinedScore": <score out of 100>,
                    "reason": "Detail reason for the combined score based on both content and delivery",
                    "suggestions": {{
                        "content": "Suggestions to improve presentation content",
                        "delivery": "Suggestions to improve delivery performance"
                    }}
                }}
            }}

            Strictly follow this return format with no additional text.
        """
        # Pass the prompt to the Groq model
        response = client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            max_tokens=2048,
            top_p=1,
            stream=False,
        )

        # Extract and parse the response
        analysis_results = response.choices[0].message.content.strip()
        # print(analysis_results)
        try:
            analysis_results = json.loads(analysis_results)
            # print(analysis_results)
        except json.JSONDecodeError:
            raise ValueError("Invalid JSON response from Groq API")

        return analysis_results
        
        # return {"message": "Training data received successfully!"}

    except Exception as error:
        print(traceback.format_exc())
        return {"Error": str(error)}

    finally:
        if os.path.exists(tmp_file_path):
            os.unlink(tmp_file_path)
    
##########################################################################################################
def generate_quiz(ppt_file_contents):
    """Generate quiz questions based on PPT content."""
    try:
        # Extract text from the PPT
        ppt_text = extract_text_from_ppt(ppt_file_contents)
        
        # Prepare the prompt for the Groq model
        prompt = f"""
        You are an AI assistant tasked with creating quiz questions based on the following presentation content:

        Content: {ppt_text}

        Create exactly 10 quiz questions based on the presentation content. Each question should:
        1. Test understanding of key concepts from the presentation
        2. Be clear and specific
        3. Have an ideal answer that is 2-3 sentences long
        4. Cover different aspects of the presentation content
        5. Be challenging but fair
        6. Do not ask questions that is out of context or not related to the presentation content
        7. Create Questions that can be asked by the presenter's audience.

        Return a JSON object with exactly this format:
        {{
            "questions": [
                {{
                    "questionNumber": 1,
                    "question": "The question text",
                    "idealAnswer": "The ideal answer in 2-3 sentences"
                }},
                ... (repeat for all 10 questions)
            ]
        }}

        Only return the JSON object with no additional text.
        """

        # Get response from Groq model
        response = client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=2048,
            top_p=1,
            stream=False,
        )

        # Extract and parse the response
        quiz_results = response.choices[0].message.content.strip()
        try:
            quiz_results = json.loads(quiz_results)
        except json.JSONDecodeError:
            raise ValueError("Invalid JSON response from Groq API")

        return quiz_results

    except Exception as error:
        print(traceback.format_exc())
        return {"Error": str(error)}
##########################################################################################################


# debug code :
#     print("DEBUG: Analyzing template structure")
#     for layout in prs.slide_layouts:
#         print(f"Layout index: {prs.slide_layouts.index(layout)}")
#         for placeholder in layout.placeholders:
#             print(f"  Placeholder index: {placeholder.placeholder_format.idx}, type: {placeholder.placeholder_format.type}")
#     print("DEBUG: Template analysis complete")

# def create_presentation(request_data):

#     template_path = f"templates/{request_data['template']}.pptx"
#     print(template_path)# Template name from request
#     if not os.path.exists(template_path):
#         raise ValueError(f"Template {request_data['template']} not found")

#     prs = Presentation(template_path) # Use the selected template
    
#     print("DEBUG: Analyzing template structure")
#     for layout in prs.slide_layouts:
#         print(f"Layout index: {prs.slide_layouts.index(layout)}")
#         for placeholder in layout.placeholders:
#             print(f"  Placeholder index: {placeholder.placeholder_format.idx}, type: {placeholder.placeholder_format.type}")
#     print("DEBUG: Template analysis complete")

#     slide_layout = prs.slide_layouts[2]  # Adjust based on template design
#     title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
#     title_slide.shapes.title.text = request_data["topic"]

#     contents_slide = prs.slides.add_slide(prs.slide_layouts[1])
#     contents_slide.shapes.title.text = "Contents"
#     slide_titles = generate_slide_titles(
#         request_data["topic"], request_data["numberOfSlides"], request_data["audienceType"], request_data["slideContent"]
#     )
#     content_text = "\n".join(slide_titles)
#     # contents_slide.shapes.placeholders[12].text = content_text
#     contents_slide.shapes.body.text = content_text
    
    
#     # Add content slides
#     for slide_title in slide_titles:
#         slide = prs.slides.add_slide(slide_layout)
#         slide.shapes.title.text = slide_title
#         slide_content = generate_slide_content(slide_title, request_data["audienceType"])
#         # slide.shapes.placeholders[12].text = slide_content
#         slide.shapes.body.text = slide_content

#     thank_you_slide = prs.slides.add_slide(prs.slide_layouts[3])
#     thank_you_slide.shapes.title.text = "Thank You"

#     ppt_stream = BytesIO()
#     prs.save(ppt_stream)
#     ppt_stream.seek(0)
#     return ppt_stream

#################################################################################################################################################

# def extract_text_from_ppt(ppt_file_contents):
#     """Extract all text from the PPT for processing."""
#     presentation = Presentation(BytesIO(ppt_file_contents))
#     all_text = []
#     for slide in presentation.slides:
#         slide_text = []
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 slide_text.append(shape.text)
#         all_text.append("\n".join(slide_text))
#     return "\n\n".join(all_text)

# def rate_ppt(ppt_file_contents):
#     """Evaluate a PPT file using the Groq model."""
#     with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
#         tmp_file.write(ppt_file_contents)
#         tmp_file_path = tmp_file.name

#     try:
#         # Extract text from the PPT for the prompt
#         ppt_text = extract_text_from_ppt(ppt_file_contents)

#         # Prepare the prompt for the Groq model
#         prompt = f"""
#         You are an AI assistant tasked with evaluating PowerPoint presentations.
        
#         PPT content: {ppt_text}
        
#         Analyze the provided presentation based on these detailed criteria of a professional presentation:
        
#         1. **Number of points per slide**: 
#             - Evaluate based on the 7/7, 8/8, or 6/6 rule (no more than 6-8 words per line and 6-8 lines per slide).
#             - Deduct points if a slide exceeds these limits.
#             - Score out of 100: Full points if all slides adhere, deductions for excess.
        
#         2. **Number of images per slide**:
#             - Find 
#             - Score out of 100: Full points for the ideal range, deductions for too few or too many images.
        
#         3. **Readability of text content**:
#             - Use the SMOG Readability Formula to assess readability (based on syllables and sentence complexity).
#             - Full points for readability suited for the target audience (e.g., 6th-9th-grade level for general audiences).
#             - Score out of 100.
        
#         4. **Consistency of slide formatting**:
#             - Check for uniformity in font styles, font sizes, color schemes, and alignment.
#             - Deduct points for inconsistent elements across slides.
#             - Score out of 100.
        
#         5. **Overall presentation content quality**:
#             - Assess the presentation's organization, logical flow, and coverage of the topic.
#             - Deduct points for missing key information or lack of structure.
#             - Score out of 100.
        
#         6. **Number of slides**:
#             - Check if the total number of slides is appropriate for the presentation's purpose (e.g., 8-12 slides for a 10-minute talk).
#             - Deduct points for excessive or insufficient slides.
#             - Score out of 100.
        
#         7. **Overall score**:
#             - A average based on all criteria above.
#             - Score out of 100.

#         Return a JSON object in this format:
#         {{
#             "noOfPoints": {{
#                 "score": <score out of 100>,
#                 "reason": "Brief reason for the score and how we computed the score"
#             }},
#             "noOfImages": {{
#                 "score": <score out of 100>,
#                 "reason": "Brief reason for the score and how we computed the score."
#             }},
#             "readability": {{
#                 "score": <score out of 100>,
#                 "reason": "Brief reason for the score and how we computed the score"
#             }},
#             "consistency": {{
#                 "score": <score out of 100>,
#                 "reason": "Brief reason for the score and how we computed the score"
#             }},
#             "quality": {{
#                 "score": <score out of 100>,
#                 "reason": "Brief reason for the score and how we computed the score"
#             }},
#             "noOfSlides": {{
#                 "score": <score out of 100>,
#                 "reason": "Brief reason for the score and how we computed the score"
#             }},
#             "overallScore": {{
#                 "score": <score out of 100>,
#                 "reason": "Brief reason for the score and how we computed the score"
#             }}
#         }}

#         Only return the JSON object with no additional text.
#         """

#         # Pass the prompt to the Groq model
#         response = client.chat.completions.create(
#             model="llama3-8b-8192",
#             messages=[
#                 {
#                     "role": "user",
#                     "content": prompt
#                 }
#             ],
#             temperature=0.7,
#             max_tokens=1024,
#             top_p=1,
#             stream=False,
#         )

#         # Extract and parse the response
#         evaluation_results = response.choices[0].message.content.strip()
#         try:
#             evaluation_results = json.loads(evaluation_results)
#         except json.JSONDecodeError:
#             raise ValueError("Invalid JSON response from Groq API")

#         return evaluation_results

#     except Exception as error:
#         print(traceback.format_exc())  # Log detailed error
#         return {"Error": str(error)}

#     finally:
#         if os.path.exists(tmp_file_path):
#             os.unlink(tmp_file_path)


####################################################################################################################
    # Save the presentation
    # os.makedirs('generated_ppt', exist_ok=True)
    # ppt_path = os.path.join('generated_ppt', f'{request_data["topic"]}_presentation.pptx')
    # print(f"////////////////////\n////////////////\n{ppt_path}")
    # prs.save(ppt_path)
    # return ppt_path



# def rate_ppt(ppt_file_contents):
#     with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
#         tmp_file.write(ppt_file_contents)
#         tmp_file_path = tmp_file.name

#     try:
#         with zipfile.ZipFile(tmp_file_path, 'r') as zip_ref:
#             # Get total number of slides
#             presentation_xml = zip_ref.read('ppt/presentation.xml')
#             root = ET.fromstring(presentation_xml)
#             total_slides = len(root.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldId'))

#             # Initialize ratings
#             slide_number_rating = 10
#             bullet_point_rating = 10

#             # Analyze each slide
#             for i in range(1, total_slides + 1):
#                 try:
#                     slide_xml = zip_ref.read(f'ppt/slides/slide{i}.xml')
#                     slide_root = ET.fromstring(slide_xml)

#                     # Check for slide number
#                     slide_number = slide_root.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ph[@type="sldNum"]')
#                     if slide_number is None:
#                         slide_number_rating -= 0.5

#                     # Count bullet points
#                     bullet_points = slide_root.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
#                     if len(bullet_points) > 7:
#                         bullet_point_rating -= 1

#                 except KeyError:
#                     # Slide doesn't exist, reduce rating
#                     slide_number_rating -= 1

#         # Ensure ratings don't go below 0
#         slide_number_rating = max(0, slide_number_rating)
#         bullet_point_rating = max(0, bullet_point_rating)

#         # Calculate overall rating
#         overall_rating = (slide_number_rating + bullet_point_rating) / 2

#         return {
#             "overall_rating": overall_rating,
#             "slide_number_rating": slide_number_rating,
#             "bullet_point_rating": bullet_point_rating,
#             "total_slides": total_slides
#         }
    
#     except Exception as error:
#     # handle the exception
#         return {"Error" : error}

#     finally:
#         # Clean up the temporary file
#         if os.path.exists(tmp_file_path):
#             os.unlink(tmp_file_path)


# def rate_ppt(ppt_file_contents):
#     """Evaluate a PPT file using the Groq model."""
#     with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
#         tmp_file.write(ppt_file_contents)
#         tmp_file_path = tmp_file.name

#     try:
#         # Extract text from the PPT for the prompt
#         ppt_text = extract_text_from_ppt(ppt_file_contents)

#         # Prepare the prompt for the Groq model
#         prompt = f"""
#         You are an AI assistant tasked with evaluating PowerPoint presentations.
        
#         PPT content: {ppt_text}
        
#         Analyze the provided presentation based on these criteria of a professional presentation:
#         1. Number of points per slide -  use 7/7 rule or 8/8 rule or 6/6 rule (score out of 100).
#         2. Number of images per slide - total number of images divided by total number of slides(score out of 100).
#         3. Readability of text content - Use SMOG Readability formula (score out of 100).
#         4. Consistency of slide formatting - (score out of 10).
#         5. Overall presentation content quality (score out of 10).
#         6.  (score out of 10).
#         7. Overall score. (score out of 10)

#         Return a JSON object in this format:
#         {{
#             "noOfPoints": {
#                 "score": ,
#                 "Reason in detail" : 
#             },
#             "noOfImages": score,
#             "Readability": score,
#             "Consistency": score,
#             "Quality": score,
#             "noOfSlides": score,
#             "overAllScore": score
#         }}

#         Only return the JSON object with no additional text.
#         """

#         # Pass the prompt to the Groq model
#         response = client.chat.completions.create(
#             model="llama3-8b-8192",
#             messages=[
#                 {
#                     "role": "user",
#                     "content": prompt
#                 }
#             ],
#             temperature=0.7,
#             max_tokens=1024,
#             top_p=1,
#             stream=False,
#         )

#         # Extract and parse the response
#         evaluation_results = response.choices[0].message.content.strip()
#         try:
#             evaluation_results = json.loads(evaluation_results)
#         except json.JSONDecodeError:
#             raise ValueError("Invalid JSON response from Groq API")

#         return evaluation_results

#     except Exception as error:
#         print(traceback.format_exc())  # Log detailed error
#         return {"Error": str(error)}

#     finally:
#         if os.path.exists(tmp_file_path):
#             os.unlink(tmp_file_path)


####################################################################################################################################333
#######################################################################################################################################
# Function to convert PPT to images
# def convert_ppt_to_images(ppt_file):
#     # Save the uploaded file to a temporary file
#     with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt_file:
#         tmp_ppt_file.write(ppt_file.getvalue())  # Write the content of the uploaded file
#         tmp_ppt_file_path = tmp_ppt_file.name    # Get the file path

#     try:
#         # Load the presentation from the saved file
#         presentation = Presentation()
#         presentation.LoadFromFile(tmp_ppt_file_path)

#         images = []
#         with tempfile.TemporaryDirectory() as tmpdirname:
#             for i, slide in enumerate(presentation.Slides):
#                 fileName = f"slide_{i}.png"
#                 image = slide.SaveAsImageByWH(800, 450)
#                 file_path = os.path.join(tmpdirname, fileName)
#                 image.Save(file_path)
#                 image.Dispose()

#                 # Open the image, convert it to RGB (to ensure compatibility), and store it in memory
#                 with Imagee.open(file_path) as img:
#                     images.append(img.copy().convert('RGB'))

#         presentation.Dispose()
#         return images

#     finally:
#         # Clean up the temporary file
#         if os.path.exists(tmp_ppt_file_path):
#             os.remove(tmp_ppt_file_path)





# Function to transcribe audio using SpeechRecognition
# def transcribe_audio(audio_file):
    # recognizer = sr.Recognizer()
    # with sr.AudioFile(audio_file) as source:
    #     audio_data = recognizer.record(source)
    #     try:
    #         text = recognizer.recognize_google(audio_data)
    #         return text
    #     except sr.UnknownValueError:
    #         return "Audio is unclear. Could not transcribe."
    #     except sr.RequestError as e:
    #         return f"Could not request results; {e}"